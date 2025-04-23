import tkinter as tk
from tkinter import filedialog, messagebox
import ttkbootstrap as ttk
from ttkbootstrap.constants import *
import json
import os
import threading
import multiprocessing
from pdf2docx import Converter
from docx2pdf import convert
from concurrent.futures import ProcessPoolExecutor, as_completed
from PIL import Image, ImageTk
import ctypes
import datetime
from pdf2image import convert_from_path
import io
import msoffcrypto
import shutil
from win32com import client
import win32com.client
import requests
import sys
import zipfile
from threading import Thread
from urllib.parse import urljoin

# 设置DPI感知
try:
    ctypes.windll.shcore.SetProcessDpiAwareness(1)
except:
    pass

def convert_file(file_path, save_path, conv_type, quality, config):
    try:
        file_name = os.path.basename(file_path)
        file_ext = os.path.splitext(file_name)[1].lower()
        
        # 根据转换类型和配置确定输出格式
        if conv_type == "pdf2word":
            output_ext = f".{config.get('word_format', 'docx')}"
        elif conv_type == "word2pdf":
            output_ext = ".pdf"
        elif conv_type == "excel2pdf":
            output_ext = ".pdf"
        elif conv_type == "pdf2excel":
            output_ext = f".{config.get('excel_format', 'xlsx')}"
        elif conv_type == "ppt2pdf":
            output_ext = ".pdf"
        elif conv_type == "pdf2ppt":
            output_ext = f".{config.get('ppt_format', 'pptx')}"
        elif conv_type == "image2pdf":
            output_ext = ".pdf"
        elif conv_type == "pdf2image":
            output_ext = f".{config.get('image_format', 'png')}"
        else:
            output_ext = ".pdf"
        
        # 构建输出文件路径
        output_path = os.path.join(save_path, os.path.splitext(file_name)[0] + output_ext)
        
        if conv_type == "pdf2word":
            cv = Converter(file_path)
            # 根据质量设置调整转换参数
            if quality == "fast":
                cv.convert(output_path, pages=None, start=0, end=None)
            elif quality == "normal":
                cv.convert(output_path, pages=None, start=0, end=None, 
                         layout_analysis=True, table_analysis=True)
            else:  # high
                cv.convert(output_path, pages=None, start=0, end=None, 
                         layout_analysis=True, table_analysis=True, 
                         image_analysis=True)
            cv.close()
        elif conv_type == "word2pdf":
            convert(file_path, output_path)
        elif conv_type == "excel2pdf":
            from win32com import client
            excel = client.Dispatch("Excel.Application")
            excel.Visible = False
            wb = excel.Workbooks.Open(file_path)
            wb.ExportAsFixedFormat(0, output_path)
            wb.Close()
            excel.Quit()
        elif conv_type == "pdf2excel":
            # 使用pdf2docx将PDF转换为Word，然后使用python-docx处理
            from pdf2docx import Converter
            from docx import Document
            import pandas as pd
            
            # 先转换为Word
            temp_docx = output_path.replace('.xlsx', '.docx')
            cv = Converter(file_path)
            cv.convert(temp_docx)
            cv.close()
            
            # 从Word中提取表格数据
            doc = Document(temp_docx)
            data = []
            for table in doc.tables:
                for row in table.rows:
                    data.append([cell.text for cell in row.cells])
            
            # 保存为Excel
            df = pd.DataFrame(data)
            df.to_excel(output_path, index=False, header=False)
            
            # 删除临时文件
            os.remove(temp_docx)
        elif conv_type == "ppt2pdf":
            from win32com import client
            powerpoint = client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = False
            presentation = powerpoint.Presentations.Open(file_path)
            presentation.SaveAs(output_path, 32)  # 32是PDF格式
            presentation.Close()
            powerpoint.Quit()
        elif conv_type == "pdf2ppt":
            # 使用pdf2docx将PDF转换为Word，然后使用python-pptx创建PPT
            from pdf2docx import Converter
            from docx import Document
            from pptx import Presentation
            
            # 先转换为Word
            temp_docx = output_path.replace('.pptx', '.docx')
            cv = Converter(file_path)
            cv.convert(temp_docx)
            cv.close()
            
            # 从Word中提取内容创建PPT
            doc = Document(temp_docx)
            prs = Presentation()
            
            for para in doc.paragraphs:
                if para.text.strip():
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    content = slide.placeholders[1]
                    title.text = para.text
                    content.text = ""  # 清空内容占位符
            
            prs.save(output_path)
            
            # 删除临时文件
            os.remove(temp_docx)
        
        return {"success": True, "file_name": file_name, "output_path": output_path}
    except Exception as e:
        return {"success": False, "file_name": file_name, "error": str(e)}

class UpdateManager:
    def __init__(self, app):
        self.app = app
        self.update_url = "https://api.github.com/repos/alove77580/fileConvert/releases/latest"  # 修改为正确的仓库地址
        self.current_version = "1.0.0"  # 当前版本号
        self.download_path = "update.zip"
        
    def check_update(self):
        """检查是否有新版本"""
        try:
            response = requests.get(self.update_url)
            if response.status_code == 200:
                latest_release = response.json()
                latest_version = latest_release['tag_name'].lstrip('v')
                
                if self.compare_versions(latest_version, self.current_version) > 0:
                    return {
                        'available': True,
                        'version': latest_version,
                        'url': latest_release['assets'][0]['browser_download_url'],
                        'description': latest_release['body']
                    }
            elif response.status_code == 404:
                print("更新检查失败：仓库不存在或无法访问")
                return {'available': False, 'error': '仓库不存在或无法访问'}
            else:
                print(f"更新检查失败：HTTP {response.status_code}")
                return {'available': False, 'error': f'HTTP {response.status_code}'}
            return {'available': False}
        except Exception as e:
            print(f"检查更新时出错: {str(e)}")
            return {'available': False, 'error': str(e)}
    
    def compare_versions(self, v1, v2):
        """比较版本号"""
        v1_parts = list(map(int, v1.split('.')))
        v2_parts = list(map(int, v2.split('.')))
        
        for i in range(max(len(v1_parts), len(v2_parts))):
            v1_part = v1_parts[i] if i < len(v1_parts) else 0
            v2_part = v2_parts[i] if i < len(v2_parts) else 0
            
            if v1_part > v2_part:
                return 1
            elif v1_part < v2_part:
                return -1
        return 0
    
    def download_update(self, url):
        """下载更新包"""
        try:
            response = requests.get(url, stream=True)
            total_size = int(response.headers.get('content-length', 0))
            
            with open(self.download_path, 'wb') as f:
                downloaded = 0
                for data in response.iter_content(chunk_size=4096):
                    downloaded += len(data)
                    f.write(data)
                    # 更新下载进度
                    if total_size > 0:
                        progress = (downloaded / total_size) * 100
                        self.app.update_download_progress(progress)
            
            return True
        except Exception as e:
            print(f"下载更新时出错: {str(e)}")
            return False
    
    def install_update(self):
        """安装更新"""
        try:
            # 解压更新包
            with zipfile.ZipFile(self.download_path, 'r') as zip_ref:
                zip_ref.extractall("temp_update")
            
            # 复制新文件
            for root, dirs, files in os.walk("temp_update"):
                for file in files:
                    src_path = os.path.join(root, file)
                    dst_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 
                                          os.path.relpath(src_path, "temp_update"))
                    os.makedirs(os.path.dirname(dst_path), exist_ok=True)
                    shutil.copy2(src_path, dst_path)
            
            # 清理临时文件
            shutil.rmtree("temp_update")
            os.remove(self.download_path)
            
            return True
        except Exception as e:
            print(f"安装更新时出错: {str(e)}")
            return False

class FileConverterApp:
    def __init__(self, root):
        self.root = root
        self.root.title("星空-文件转换工具")
        self.root.geometry("1000x1500")  # 设置初始大小
        self.root.minsize(1000, 1500)  # 设置最小大小
        self.root.resizable(True, True)  # 允许调整大小
        
        # 初始化配置
        self.config = self.load_config()
        
        # 初始化自动打开变量
        self.auto_open_var = tk.BooleanVar(value=self.config.get('auto_open', False))
        
        # 初始化更新管理器
        self.update_manager = UpdateManager(self)
        
        # 创建更新进度窗口
        self.create_update_window()
        
        # 创建历史记录窗口
        self.history_window = tk.Toplevel(root)
        self.history_window.title("转换历史")
        self.history_window.geometry("800x600")  # 设置初始大小
        self.history_window.minsize(600, 400)  # 设置最小大小
        self.history_window.resizable(True, True)  # 允许调整大小
        self.history_window_visible = True
        
        # 创建历史记录主框架
        self.history_main_frame = ttk.Frame(self.history_window, padding="20")
        self.history_main_frame.pack(fill=tk.BOTH, expand=True)
        
        # 创建历史记录标题栏
        self.history_title_frame = ttk.Frame(self.history_main_frame)
        self.history_title_frame.pack(fill=tk.X, pady=(0, 10))
        
        # 创建历史记录标题
        self.history_header = ttk.Label(
            self.history_title_frame,
            text="转换历史记录",
            font=('微软雅黑', 16, 'bold'),
            bootstyle="light"
        )
        self.history_header.pack(side=tk.LEFT, expand=True)
        
        # 添加关闭按钮
        self.close_history_button = ttk.Button(
            self.history_title_frame,
            text="×",
            command=self.toggle_history_window,
            width=3,
            bootstyle="danger"
        )
        self.close_history_button.pack(side=tk.RIGHT)
        
        # 创建历史记录列表框架
        list_frame = ttk.Frame(self.history_main_frame)
        list_frame.pack(fill=tk.BOTH, expand=True, pady=5)
        
        # 创建水平滚动条
        h_scrollbar = ttk.Scrollbar(
            list_frame,
            orient=tk.HORIZONTAL,
            bootstyle="round"
        )
        h_scrollbar.pack(side=tk.BOTTOM, fill=tk.X)
        
        # 创建垂直滚动条
        v_scrollbar = ttk.Scrollbar(
            list_frame,
            orient=tk.VERTICAL,
            bootstyle="round"
        )
        v_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # 创建历史记录列表
        self.history_list = tk.Listbox(
            list_frame,
            width=50,
            height=20,
            font=('微软雅黑', 10),
            bg='#2b3e50',
            fg='white',
            selectbackground='#0078d7',
            selectforeground='white',
            relief=tk.FLAT,
            xscrollcommand=h_scrollbar.set,
            yscrollcommand=v_scrollbar.set
        )
        self.history_list.pack(fill=tk.BOTH, expand=True)
        
        # 配置滚动条
        h_scrollbar.config(command=self.history_list.xview)
        v_scrollbar.config(command=self.history_list.yview)
        
        # 创建底部按钮框架
        bottom_frame = ttk.Frame(self.history_main_frame)
        bottom_frame.pack(fill=tk.X, pady=(10, 0), side=tk.BOTTOM)
        
        # 重新执行按钮
        self.reconvert_button = ttk.Button(
            bottom_frame,
            text="重新执行",
            command=self.reconvert_history,
            width=15,
            bootstyle="info"
        )
        self.reconvert_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 清空历史按钮
        self.clear_history_button = ttk.Button(
            bottom_frame,
            text="清空历史",
            command=self.clear_history,
            width=15,
            bootstyle="danger"
        )
        self.clear_history_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 加载历史记录
        self.load_history()
        
        # 创建主框架
        self.main_frame = ttk.Frame(root, padding="20")
        self.main_frame.pack(fill=tk.BOTH, expand=True)
        
        # 创建标题区域
        self.header_frame = ttk.Frame(self.main_frame)
        self.header_frame.pack(fill=tk.X, pady=(0, 10))
        
        # 初始化并行处理数量
        self.parallel_count = tk.StringVar(value="4")  # 默认使用4个进程
        
        # 创建质量选项的中文映射
        self.quality_map = {
            "high": "高质量",
            "normal": "标准",
            "low": "低质量"
        }
        self.quality_reverse_map = {
            "高质量": "high",
            "标准": "normal",
            "低质量": "low"
        }
        
        # 设置主题
        self.style = ttk.Style()
        theme = self.config.get('theme', 'superhero')
        self.style.theme_use(theme)
        
        # 初始化主题颜色
        self.update_all_windows_theme()
        
        # 应用自定义颜色
        self.style.configure('.', 
                           foreground=self.config.get('fg_color', '#ffffff'),
                           background=self.config.get('bg_color', '#2b3e50'))
        
        # 创建标题
        self.header_label = ttk.Label(
            self.header_frame, 
            text="文件转换工具", 
            font=('微软雅黑', 24, 'bold'),
            bootstyle="light"
        )
        self.header_label.pack(expand=True)
        
        # 添加高级设置按钮
        self.advanced_settings_button = ttk.Button(
            self.header_frame,
            text="高级设置",
            command=self.show_advanced_settings,
            width=10,
            bootstyle="info"
        )
        self.advanced_settings_button.pack(side=tk.RIGHT, padx=10)
        
        # 添加更新按钮（放在高级设置按钮旁边）
        self.update_button = ttk.Button(
            self.header_frame,
            text="检查更新",
            command=self.check_for_updates,
            width=10,
            bootstyle="info"
        )
        self.update_button.pack(side=tk.RIGHT, padx=5)
        
        # 添加历史记录控制按钮
        self.history_control_button = ttk.Button(
            self.header_frame,
            text=">>",
            command=self.toggle_history_window,
            width=3,
            bootstyle="info"
        )
        self.history_control_button.pack(side=tk.RIGHT, padx=5)
        
        # 绑定窗口关闭事件
        self.root.protocol("WM_DELETE_WINDOW", self.on_closing)
        self.history_window.protocol("WM_DELETE_WINDOW", self.on_closing)
        
        # 绑定快捷键
        self.bind_shortcuts()
        
        # 设置窗口居中显示
        self.root.update_idletasks()
        width = self.root.winfo_width()
        height = self.root.winfo_height()
        x = (self.root.winfo_screenwidth() // 2) - (width // 2)
        y = (self.root.winfo_screenheight() // 2) - (height // 2)
        self.root.geometry(f'{width}x{height}+{x}+{y}')
        
        # 设置历史记录窗口位置（在主窗口右侧）
        self.history_window.update_idletasks()
        history_x = x + width + 10  # 在主窗口右侧10像素
        self.history_window.geometry(f'{width}x{height}+{history_x}+{y}')
        
        # 绑定主窗口移动事件
        self.root.bind('<Configure>', self.on_root_move)
        
        # 初始化变量
        self.files_to_convert = []
        self.conversion_type = tk.StringVar(value="auto")
        
        # 创建文件选择区域
        self.file_frame = ttk.Labelframe(
            self.main_frame, 
            text="文件选择", 
            padding="15",
            bootstyle="info"
        )
        self.file_frame.pack(fill=tk.BOTH, expand=True, pady=5, padx=50)
        
        # 创建文件列表显示区域
        self.list_frame = ttk.Frame(self.file_frame)
        self.list_frame.pack(fill=tk.BOTH, expand=True, pady=5)
        
        # 创建文件列表
        self.file_list = tk.Listbox(
            self.list_frame, 
            width=50, 
            height=8,
            font=('微软雅黑', 10),
            bg='#2b3e50',
            fg='white',
            selectbackground='#0078d7',
            selectforeground='white',
            relief=tk.FLAT
        )
        self.file_list.pack(fill=tk.BOTH, expand=True, side=tk.LEFT)
        
        # 添加滚动条
        scrollbar = ttk.Scrollbar(
            self.list_frame, 
            orient=tk.VERTICAL, 
            command=self.file_list.yview,
            bootstyle="round"
        )
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.file_list.config(yscrollcommand=scrollbar.set)
        
        # 创建按钮区域
        self.button_frame = ttk.Frame(self.file_frame)
        self.button_frame.pack(fill=tk.X, pady=5)
        
        self.select_button = ttk.Button(
            self.button_frame, 
            text="选择文件", 
            command=self.select_files,
            width=15,
            bootstyle="primary"
        )
        self.select_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        self.convert_button = ttk.Button(
            self.button_frame, 
            text="开始转换", 
            command=self.start_conversion,
            width=15,
            bootstyle="success"
        )
        self.convert_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 在文件列表区域添加操作按钮
        self.list_buttons_frame = ttk.Frame(self.file_frame)
        self.list_buttons_frame.pack(fill=tk.X, pady=(0, 5))
        
        # 添加排序按钮
        self.sort_button = ttk.Button(
            self.list_buttons_frame,
            text="排序",
            command=self.show_sort_dialog,
            width=10,
            bootstyle="info"
        )
        self.sort_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 添加合并按钮
        self.merge_button = ttk.Button(
            self.list_buttons_frame,
            text="合并",
            command=self.show_merge_dialog,
            width=10,
            bootstyle="info"
        )
        self.merge_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 添加分割按钮
        self.split_button = ttk.Button(
            self.list_buttons_frame,
            text="分割",
            command=self.show_split_dialog,
            width=10,
            bootstyle="info"
        )
        self.split_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 添加加密/解密按钮
        self.encrypt_button = ttk.Button(
            self.list_buttons_frame,
            text="加密/解密",
            command=self.show_encrypt_dialog,
            width=10,
            bootstyle="info"
        )
        self.encrypt_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        self.clear_button = ttk.Button(
            self.list_buttons_frame,
            text="清空列表",
            command=self.clear_file_list,
            width=10,
            bootstyle="danger"
        )
        self.clear_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        self.remove_button = ttk.Button(
            self.list_buttons_frame,
            text="删除选中",
            command=self.remove_selected_files,
            width=10,
            bootstyle="warning"
        )
        self.remove_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        self.rename_button = ttk.Button(
            self.list_buttons_frame,
            text="批量重命名",
            command=self.show_rename_dialog,
            width=10,
            bootstyle="info"
        )
        self.rename_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 创建文件预览区域
        self.preview_frame = ttk.Labelframe(
            self.file_frame, 
            text="文件预览", 
            padding="10",
            bootstyle="info"
        )
        self.preview_frame.pack(fill=tk.BOTH, expand=True, pady=5)
        
        # 创建预览内容区域
        self.preview_content = ttk.Frame(self.preview_frame)
        self.preview_content.pack(fill=tk.BOTH, expand=True)
        
        # 文件信息标签
        self.file_info_label = ttk.Label(
            self.preview_content,
            text="选择文件后显示预览信息",
            font=('微软雅黑', 10),
            bootstyle="light"
        )
        self.file_info_label.pack(pady=5)
        
        # PDF预览图像
        self.preview_image = None
        self.preview_image_label = ttk.Label(
            self.preview_content,
            bootstyle="light"
        )
        self.preview_image_label.pack(pady=5)
        
        # 绑定文件列表选择事件
        self.file_list.bind('<<ListboxSelect>>', self.on_file_select)
        
        # 创建进度条区域
        self.progress_frame = ttk.Labelframe(
            self.main_frame, 
            text="转换进度", 
            padding="20",
            bootstyle="info"
        )
        self.progress_frame.pack(fill=tk.X, pady=10, padx=100, expand=True)
        
        # 创建进度条容器
        progress_container = ttk.Frame(self.progress_frame)
        progress_container.pack(fill=tk.BOTH, expand=True, pady=5)
        
        self.progress = ttk.Progressbar(
            progress_container, 
            length=400, 
            mode='determinate',
            bootstyle="success-striped"
        )
        self.progress.pack(fill=tk.X, expand=True, pady=5)
        
        # 创建状态标签
        self.status_label = ttk.Label(
            progress_container, 
            text="就绪", 
            font=('微软雅黑', 10),
            bootstyle="light"
        )
        self.status_label.pack(pady=5)
        
        # 添加统计信息标签
        self.stats_label = ttk.Label(
            progress_container,
            text="",
            font=('微软雅黑', 9),
            bootstyle="light"
        )
        self.stats_label.pack(pady=5)
        
        # 初始化统计变量
        self.total_converted = 0
        self.successful_conversions = 0
        self.failed_conversions = 0
        
        # 初始化模板数据
        self.templates = self.load_templates()
        
        # 在文件选择区域添加模板按钮
        self.template_button = ttk.Button(
            self.button_frame,
            text="模板管理",
            command=self.show_template_dialog,
            width=15,
            bootstyle="info"
        )
        self.template_button.pack(side=tk.LEFT, padx=5, expand=True)
    
    def load_config(self):
        """加载配置文件"""
        try:
            with open('config.json', 'r', encoding='utf-8') as f:
                config = json.load(f)
                # 确保所有必要的配置项都存在
                if 'pdf_quality' not in config:
                    config['pdf_quality'] = 'high'
                if 'auto_open' not in config:
                    config['auto_open'] = True
                if 'word_format' not in config:
                    config['word_format'] = 'docx'
                if 'excel_format' not in config:
                    config['excel_format'] = 'xlsx'
                if 'ppt_format' not in config:
                    config['ppt_format'] = 'pptx'
                if 'pdf_dpi' not in config:
                    config['pdf_dpi'] = 300
                if 'theme' not in config:
                    config['theme'] = 'superhero'
                return config
        except FileNotFoundError:
            # 如果配置文件不存在，创建默认配置
            default_config = {
                'pdf_quality': 'high',
                'auto_open': True,
                'word_format': 'docx',
                'excel_format': 'xlsx',
                'ppt_format': 'pptx',
                'pdf_dpi': 300,
                'theme': 'superhero'
            }
            self.save_config(default_config)
            return default_config
        except Exception as e:
            messagebox.showerror("错误", f"加载配置文件时出错：{str(e)}")
            return self.load_config()  # 出错时返回默认配置

    def save_config(self, config=None):
        """保存配置文件"""
        try:
            with open('config.json', 'w', encoding='utf-8') as f:
                json.dump(config or self.config, f, ensure_ascii=False, indent=4)
        except Exception as e:
            messagebox.showerror("错误", f"保存配置文件时出错：{str(e)}")
    
    def load_history(self):
        """加载历史记录"""
        try:
            with open('history.json', 'r', encoding='utf-8') as f:
                self.history = json.load(f)
                # 检查并转换旧格式记录
                for record in self.history:
                    if 'operation' not in record:
                        # 旧格式记录默认为转换操作
                        record['operation'] = 'convert'
                    if 'success' not in record:
                        record['success'] = True
                    if 'timestamp' not in record:
                        record['timestamp'] = datetime.datetime.now().timestamp()
        except FileNotFoundError:
            self.history = []
        except Exception as e:
            messagebox.showerror("错误", f"加载历史记录时出错：{str(e)}")
            self.history = []
        
        self.update_history_list()
    
    def save_history(self):
        """保存历史记录"""
        try:
            with open('history.json', 'w', encoding='utf-8') as f:
                json.dump(self.history, f, ensure_ascii=False, indent=4)
        except Exception as e:
            messagebox.showerror("错误", f"保存历史记录时出错：{str(e)}")
    
    def update_history_list(self):
        """更新历史记录列表显示"""
        self.history_list.delete(0, tk.END)
        for record in self.history:
            try:
                timestamp = datetime.datetime.fromtimestamp(record['timestamp']).strftime('%Y-%m-%d %H:%M:%S')
                operation = record['operation']
                file_name = os.path.basename(record['file_path'])
                status = "成功" if record['success'] else "失败"
                
                if operation == 'convert':
                    target_format = record.get('target_format', '')
                    display_text = f"{timestamp} - 转换 {file_name} 到 {target_format} - {status}"
                elif operation == 'split':
                    display_text = f"{timestamp} - 分割 {file_name} - {status}"
                elif operation == 'merge':
                    display_text = f"{timestamp} - 合并 {len(record['file_paths'])} 个文件 - {status}"
                elif operation == 'encrypt':
                    display_text = f"{timestamp} - 加密 {file_name} - {status}"
                elif operation == 'decrypt':
                    display_text = f"{timestamp} - 解密 {file_name} - {status}"
                else:
                    display_text = f"{timestamp} - {operation} {file_name} - {status}"
                
                self.history_list.insert(tk.END, display_text)
            except Exception as e:
                # 跳过格式错误的记录
                continue
    
    def add_to_history(self, operation, file_path, success, output_path=None, error=None, **kwargs):
        """添加记录到历史"""
        record = {
            'timestamp': datetime.datetime.now().timestamp(),
            'operation': operation,
            'file_path': file_path,
            'success': success,
            'output_path': output_path,
            'error': error
        }
        record.update(kwargs)
        self.history.insert(0, record)  # 添加到开头
        if len(self.history) > 100:  # 限制历史记录数量
            self.history = self.history[:100]
        self.save_history()
        self.update_history_list()
    
    def start_conversion(self):
        if not self.files_to_convert:
            messagebox.showwarning("警告", "请先选择要转换的文件！")
            return
        
        # 检查文件类型是否匹配
        for file in self.files_to_convert:
            conv_type = self.get_conversion_type(file)
            if conv_type is None:
                messagebox.showerror("错误", f"文件 {os.path.basename(file)} 不是PDF或Word文件！")
                return
        
        save_path = filedialog.askdirectory()
        if not save_path:
            return
        
        self.config["default_save_path"] = save_path
        self.save_config()
        
        # 初始化统计变量
        self.total_converted = len(self.files_to_convert)
        self.successful_conversions = 0
        self.failed_conversions = 0
        
        # 重置进度条
        self.progress['value'] = 0
        self.status_label.config(text="正在转换...")
        self.update_stats()
        
        # 禁用转换按钮
        self.convert_button.config(state='disabled')
        
        # 开始转换
        thread = threading.Thread(target=self.convert_files, args=(save_path,))
        thread.start()
    
    def convert_files(self, save_path):
        processed_files = 0
        
        # 使用进程池进行并行处理
        with ProcessPoolExecutor(max_workers=int(self.parallel_count.get())) as executor:
            # 提交所有转换任务
            futures = []
            for file_path in self.files_to_convert:
                future = executor.submit(
                    convert_file, 
                    file_path, 
                    save_path, 
                    self.get_conversion_type(file_path),
                    self.config.get('pdf_quality', 'high'),
                    self.config
                )
                futures.append(future)
            
            # 处理完成的任务
            for future in as_completed(futures):
                result = future.result()
                processed_files += 1
                
                if result["success"]:
                    self.successful_conversions += 1
                    if self.auto_open_var.get():
                        os.startfile(result["output_path"])
                    # 添加到历史记录
                    self.add_to_history(
                        'convert',
                        result["file_name"],
                        True,
                        result["output_path"],
                        target_format=os.path.splitext(result["output_path"])[1][1:]
                    )
                else:
                    self.failed_conversions += 1
                    # 添加到历史记录
                    self.add_to_history(
                        'convert',
                        result["file_name"],
                        False,
                        error=result["error"]
                    )
                
                # 更新进度条
                progress_value = (processed_files / self.total_converted) * 100
                self.progress['value'] = progress_value
                self.progress.update()
                
                # 更新统计信息
                self.update_stats()
                
                # 更新状态标签
                if result["success"]:
                    self.status_label.config(text=f"已完成: {result['file_name']}")
                else:
                    messagebox.showerror("错误", f"转换 {result['file_name']} 时出错: {result['error']}")
                
                # 强制更新界面
                self.root.update_idletasks()
        
        # 转换完成
        self.convert_button.config(state='normal')
        messagebox.showinfo("完成", f"转换完成！\n成功: {self.successful_conversions} 个\n失败: {self.failed_conversions} 个")
        
        # 清空文件列表
        self.files_to_convert = []
        self.update_file_list()
    
    def update_stats(self):
        stats_text = f"总文件: {self.total_converted} | "
        stats_text += f"成功: {self.successful_conversions} | "
        stats_text += f"失败: {self.failed_conversions}"
        self.stats_label.config(text=stats_text)
    
    def clear_file_list(self):
        if self.files_to_convert:
            if messagebox.askyesno("确认", "确定要清空文件列表吗？"):
                self.files_to_convert = []
                self.update_file_list()
                self.status_label.config(text="文件列表已清空")
    
    def remove_selected_files(self):
        selected_indices = self.file_list.curselection()
        if selected_indices:
            # 从后往前删除，避免索引变化
            for index in reversed(selected_indices):
                del self.files_to_convert[index]
            self.update_file_list()
            self.status_label.config(text=f"已删除 {len(selected_indices)} 个文件，剩余 {len(self.files_to_convert)} 个文件")
    
    def on_file_select(self, event):
        selection = self.file_list.curselection()
        if not selection:
            return
            
        file_path = self.files_to_convert[selection[0]]
        file_name = os.path.basename(file_path)
        
        # 获取文件信息
        file_size = os.path.getsize(file_path) / 1024  # KB
        create_time = datetime.datetime.fromtimestamp(os.path.getctime(file_path))
        
        # 更新文件信息
        info_text = f"文件名: {file_name}\n"
        info_text += f"大小: {file_size:.2f} KB\n"
        info_text += f"创建时间: {create_time.strftime('%Y-%m-%d %H:%M:%S')}"
        self.file_info_label.config(text=info_text)
        
        # 如果是PDF文件，显示预览
        if file_name.lower().endswith('.pdf'):
            try:
                # 转换PDF第一页为图像
                images = convert_from_path(file_path, first_page=1, last_page=1)
                if images:
                    # 调整图像大小
                    image = images[0]
                    max_size = (300, 400)
                    image.thumbnail(max_size, Image.Resampling.LANCZOS)
                    
                    # 转换为Tkinter可用的图像
                    photo = ImageTk.PhotoImage(image)
                    self.preview_image = photo  # 保持引用
                    self.preview_image_label.config(image=photo)
            except Exception as e:
                self.preview_image_label.config(text=f"无法预览PDF: {str(e)}")
        else:
            self.preview_image_label.config(image='')
            self.preview_image_label.config(text="")
    
    def show_rename_dialog(self):
        if not self.files_to_convert:
            messagebox.showwarning("警告", "请先选择要重命名的文件！")
            return
            
        # 创建重命名对话框
        dialog = tk.Toplevel(self.root)
        dialog.title("批量重命名")
        dialog.geometry("500x480")
        dialog.resizable(False, False)
        dialog.transient(self.root)
        dialog.grab_set()
        
        # 设置对话框居中显示
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f'{width}x{height}+{x}+{y}')
        
        # 创建重命名规则框架
        rule_frame = ttk.Frame(dialog, padding="10")
        rule_frame.pack(fill=tk.BOTH, expand=True)
        
        # 前缀
        ttk.Label(
            rule_frame,
            text="前缀:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.prefix_var = tk.StringVar()
        ttk.Entry(
            rule_frame,
            textvariable=self.prefix_var,
            width=30
        ).pack(fill=tk.X, pady=(0, 10))
        
        # 后缀
        ttk.Label(
            rule_frame,
            text="后缀:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.suffix_var = tk.StringVar()
        ttk.Entry(
            rule_frame,
            textvariable=self.suffix_var,
            width=30
        ).pack(fill=tk.X, pady=(0, 10))
        
        # 起始序号
        ttk.Label(
            rule_frame,
            text="起始序号:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.start_num_var = tk.StringVar(value="1")
        ttk.Entry(
            rule_frame,
            textvariable=self.start_num_var,
            width=10
        ).pack(anchor=tk.W, pady=(0, 10))
        
        # 序号位数
        ttk.Label(
            rule_frame,
            text="序号位数:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.num_digits_var = tk.StringVar(value="3")
        ttk.Entry(
            rule_frame,
            textvariable=self.num_digits_var,
            width=10
        ).pack(anchor=tk.W, pady=(0, 10))
        
        # 创建按钮框架
        button_frame = ttk.Frame(rule_frame)
        button_frame.pack(fill=tk.X, pady=10)
        
        # 预览按钮
        preview_button = ttk.Button(
            button_frame,
            text="预览",
            command=lambda: self.preview_rename(dialog),
            width=10,
            bootstyle="info"
        )
        preview_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 应用按钮
        apply_button = ttk.Button(
            button_frame,
            text="应用",
            command=lambda: self.apply_rename(dialog),
            width=10,
            bootstyle="success"
        )
        apply_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 取消按钮
        cancel_button = ttk.Button(
            button_frame,
            text="取消",
            command=dialog.destroy,
            width=10,
            bootstyle="danger"
        )
        cancel_button.pack(side=tk.LEFT, padx=5, expand=True)
    
    def preview_rename(self, dialog):
        try:
            start_num = int(self.start_num_var.get())
            num_digits = int(self.num_digits_var.get())
        except ValueError:
            messagebox.showerror("错误", "序号和位数必须是数字！")
            return
            
        preview_text = "重命名预览：\n\n"
        for i, file_path in enumerate(self.files_to_convert):
            file_name = os.path.basename(file_path)
            file_ext = os.path.splitext(file_name)[1]
            
            # 生成序号
            num = start_num + i
            num_str = str(num).zfill(num_digits)
            
            # 生成新文件名
            new_name = f"{self.prefix_var.get()}{num_str}{self.suffix_var.get()}{file_ext}"
            preview_text += f"{file_name} -> {new_name}\n"
            
        # 显示预览
        preview_dialog = tk.Toplevel(dialog)
        preview_dialog.title("重命名预览")
        preview_dialog.geometry("500x800")
        preview_dialog.resizable(False, False)
        preview_dialog.transient(dialog)
        
        # 创建预览文本区域
        preview_text_widget = tk.Text(
            preview_dialog,
            wrap=tk.WORD,
            font=('微软雅黑', 10)
        )
        preview_text_widget.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        preview_text_widget.insert(tk.END, preview_text)
        preview_text_widget.config(state='disabled')
        
        # 关闭按钮
        close_button = ttk.Button(
            preview_dialog,
            text="关闭",
            command=preview_dialog.destroy,
            bootstyle="primary"
        )
        close_button.pack(pady=10)
    
    def apply_rename(self, dialog):
        try:
            start_num = int(self.start_num_var.get())
            num_digits = int(self.num_digits_var.get())
        except ValueError:
            messagebox.showerror("错误", "序号和位数必须是数字！")
            return
            
        # 创建重命名后的文件列表
        renamed_files = []
        for i, file_path in enumerate(self.files_to_convert):
            file_dir = os.path.dirname(file_path)
            file_name = os.path.basename(file_path)
            file_ext = os.path.splitext(file_name)[1]
            
            # 生成序号
            num = start_num + i
            num_str = str(num).zfill(num_digits)
            
            # 生成新文件名
            new_name = f"{self.prefix_var.get()}{num_str}{self.suffix_var.get()}{file_ext}"
            new_path = os.path.join(file_dir, new_name)
            
            # 检查新文件名是否已存在
            if os.path.exists(new_path):
                messagebox.showerror("错误", f"文件 {new_name} 已存在！")
                return
                
            renamed_files.append((file_path, new_path))
        
        # 执行重命名
        for old_path, new_path in renamed_files:
            try:
                os.rename(old_path, new_path)
            except Exception as e:
                messagebox.showerror("错误", f"重命名 {os.path.basename(old_path)} 时出错: {str(e)}")
                return
        
        # 更新文件列表
        self.files_to_convert = [new_path for _, new_path in renamed_files]
        self.update_file_list()
        
        # 关闭对话框
        dialog.destroy()
        messagebox.showinfo("成功", "文件重命名完成！")
    
    def on_closing(self):
        """处理窗口关闭事件"""
        # 保存配置和历史记录
        self.save_config()
        self.save_history()
        # 关闭两个窗口
        self.root.destroy()  # 只销毁主窗口，历史窗口会随之关闭

    def toggle_history_window(self):
        """切换历史记录窗口的显示状态"""
        if self.history_window_visible:
            self.history_window.withdraw()  # 隐藏窗口
            self.history_control_button.config(text="<<")
        else:
            self.history_window.deiconify()  # 显示窗口
            self.history_control_button.config(text=">>")
            # 更新窗口位置
            self.on_root_move(None)
        self.history_window_visible = not self.history_window_visible

    def on_root_move(self, event):
        """处理主窗口移动事件"""
        if event is None or event.widget == self.root:  # 确保是主窗口的移动事件或无事件
            # 获取主窗口当前位置
            root_x = self.root.winfo_x()
            root_y = self.root.winfo_y()
            root_width = self.root.winfo_width()
            root_height = self.root.winfo_height()
            
            # 更新历史记录窗口位置和大小
            history_x = root_x + root_width + 10  # 在主窗口右侧10像素
            self.history_window.geometry(f'{root_width}x{root_height}+{history_x}+{root_y}')

    def on_window_minimize(self, event):
        """处理主窗口最小化事件"""
        if self.history_window_visible:
            self.history_window.withdraw()

    def on_window_restore(self, event):
        """处理主窗口恢复事件"""
        if self.history_window_visible:
            self.history_window.deiconify()
            # 更新窗口位置
            self.on_root_move(None)

    def show_sort_dialog(self):
        """显示排序对话框"""
        if not self.files_to_convert:
            messagebox.showwarning("警告", "请先选择要排序的文件！")
            return
            
        # 创建排序对话框
        dialog = tk.Toplevel(self.root)
        dialog.title("文件排序")
        dialog.geometry("400x530")
        dialog.resizable(False, False)
        dialog.transient(self.root)
        dialog.grab_set()
        
        # 设置对话框居中显示
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f'{width}x{height}+{x}+{y}')
        
        # 创建排序选项框架
        sort_frame = ttk.Frame(dialog, padding="20")
        sort_frame.pack(fill=tk.BOTH, expand=True)
        
        # 排序方式
        ttk.Label(
            sort_frame,
            text="排序方式:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 10))
        
        self.sort_type_var = tk.StringVar(value="name")
        ttk.Radiobutton(
            sort_frame,
            text="按文件名",
            variable=self.sort_type_var,
            value="name",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=5)
        
        ttk.Radiobutton(
            sort_frame,
            text="按文件大小",
            variable=self.sort_type_var,
            value="size",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=5)
        
        ttk.Radiobutton(
            sort_frame,
            text="按创建日期",
            variable=self.sort_type_var,
            value="date",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=5)
        
        # 排序顺序
        ttk.Label(
            sort_frame,
            text="排序顺序:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(20, 10))
        
        self.sort_order_var = tk.StringVar(value="asc")
        ttk.Radiobutton(
            sort_frame,
            text="升序",
            variable=self.sort_order_var,
            value="asc",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=5)
        
        ttk.Radiobutton(
            sort_frame,
            text="降序",
            variable=self.sort_order_var,
            value="desc",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=5)
        
        # 创建按钮框架
        button_frame = ttk.Frame(sort_frame)
        button_frame.pack(fill=tk.X, pady=20)
        
        # 应用按钮
        apply_button = ttk.Button(
            button_frame,
            text="应用",
            command=lambda: self.apply_sort(dialog),
            width=10,
            bootstyle="success"
        )
        apply_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 取消按钮
        cancel_button = ttk.Button(
            button_frame,
            text="取消",
            command=dialog.destroy,
            width=10,
            bootstyle="danger"
        )
        cancel_button.pack(side=tk.LEFT, padx=5, expand=True)
    
    def apply_sort(self, dialog):
        """应用排序"""
        sort_type = self.sort_type_var.get()
        sort_order = self.sort_order_var.get()
        
        # 获取文件信息
        file_info = []
        for file_path in self.files_to_convert:
            file_name = os.path.basename(file_path)
            file_size = os.path.getsize(file_path)
            create_time = os.path.getctime(file_path)
            file_info.append({
                'path': file_path,
                'name': file_name,
                'size': file_size,
                'date': create_time
            })
        
        # 根据排序类型和顺序进行排序
        if sort_type == "name":
            file_info.sort(key=lambda x: x['name'], reverse=(sort_order == "desc"))
        elif sort_type == "size":
            file_info.sort(key=lambda x: x['size'], reverse=(sort_order == "desc"))
        else:  # date
            file_info.sort(key=lambda x: x['date'], reverse=(sort_order == "desc"))
        
        # 更新文件列表
        self.files_to_convert = [info['path'] for info in file_info]
        self.update_file_list()
        
        # 关闭对话框
        dialog.destroy()
        messagebox.showinfo("成功", "文件排序完成！")

    def show_merge_dialog(self):
        """显示合并对话框"""
        if not self.files_to_convert:
            messagebox.showwarning("警告", "请先选择要合并的文件！")
            return
            
        # 检查文件类型是否一致
        file_types = set()
        for file in self.files_to_convert:
            if file.lower().endswith('.pdf'):
                file_types.add('pdf')
            elif file.lower().endswith('.docx'):
                file_types.add('docx')
        
        if len(file_types) > 1:
            messagebox.showerror("错误", "不能同时合并PDF和Word文件！")
            return
            
        if not file_types:
            messagebox.showerror("错误", "没有可合并的文件！")
            return
            
        # 创建合并对话框
        dialog = tk.Toplevel(self.root)
        dialog.title("文件合并")
        dialog.geometry("500x650")
        dialog.resizable(False, False)
        dialog.transient(self.root)
        dialog.grab_set()
        
        # 设置对话框居中显示
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f'{width}x{height}+{x}+{y}')
        
        # 创建合并选项框架
        merge_frame = ttk.Frame(dialog, padding="20")
        merge_frame.pack(fill=tk.BOTH, expand=True)
        
        # 文件类型标签
        file_type = "PDF" if 'pdf' in file_types else "Word"
        ttk.Label(
            merge_frame,
            text=f"合并 {file_type} 文件",
            font=('微软雅黑', 12, 'bold'),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 20))
        
        # 输出文件名
        ttk.Label(
            merge_frame,
            text="输出文件名:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.output_name_var = tk.StringVar(value=f"合并后的{file_type}文件")
        ttk.Entry(
            merge_frame,
            textvariable=self.output_name_var,
            width=40
        ).pack(fill=tk.X, pady=(0, 20))
        
        # 文件顺序
        ttk.Label(
            merge_frame,
            text="文件顺序:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 10))
        
        # 创建文件顺序列表
        self.order_list = tk.Listbox(
            merge_frame,
            width=50,
            height=8,
            font=('微软雅黑', 10),
            bg='#2b3e50',
            fg='white',
            selectbackground='#0078d7',
            selectforeground='white',
            relief=tk.FLAT
        )
        self.order_list.pack(fill=tk.BOTH, expand=True, pady=5)
        
        # 添加滚动条
        scrollbar = ttk.Scrollbar(
            merge_frame,
            orient=tk.VERTICAL,
            command=self.order_list.yview,
            bootstyle="round"
        )
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.order_list.config(yscrollcommand=scrollbar.set)
        
        # 添加文件到顺序列表
        for file in self.files_to_convert:
            self.order_list.insert(tk.END, os.path.basename(file))
        
        # 创建按钮框架
        button_frame = ttk.Frame(merge_frame)
        button_frame.pack(fill=tk.X, pady=20)
        
        # 上移按钮
        up_button = ttk.Button(
            button_frame,
            text="上移",
            command=self.move_up,
            width=10,
            bootstyle="info"
        )
        up_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 下移按钮
        down_button = ttk.Button(
            button_frame,
            text="下移",
            command=self.move_down,
            width=10,
            bootstyle="info"
        )
        down_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 应用按钮
        apply_button = ttk.Button(
            button_frame,
            text="合并",
            command=lambda: self.apply_merge(dialog, file_type),
            width=10,
            bootstyle="success"
        )
        apply_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 取消按钮
        cancel_button = ttk.Button(
            button_frame,
            text="取消",
            command=dialog.destroy,
            width=10,
            bootstyle="danger"
        )
        cancel_button.pack(side=tk.LEFT, padx=5, expand=True)
    
    def move_up(self):
        """上移选中的文件"""
        selection = self.order_list.curselection()
        if not selection or selection[0] == 0:
            return
            
        index = selection[0]
        item = self.order_list.get(index)
        self.order_list.delete(index)
        self.order_list.insert(index - 1, item)
        self.order_list.selection_set(index - 1)
    
    def move_down(self):
        """下移选中的文件"""
        selection = self.order_list.curselection()
        if not selection or selection[0] == self.order_list.size() - 1:
            return
            
        index = selection[0]
        item = self.order_list.get(index)
        self.order_list.delete(index)
        self.order_list.insert(index + 1, item)
        self.order_list.selection_set(index + 1)
    
    def apply_merge(self, dialog, file_type):
        """应用合并"""
        # 获取输出文件名
        output_name = self.output_name_var.get().strip()
        if not output_name:
            messagebox.showerror("错误", "请输入输出文件名！")
            return
            
        # 添加文件扩展名
        if not output_name.lower().endswith(f'.{file_type.lower()}'):
            output_name += f'.{file_type.lower()}'
            
        # 选择保存位置
        save_path = filedialog.asksaveasfilename(
            defaultextension=f'.{file_type.lower()}',
            initialfile=output_name,
            filetypes=[(f'{file_type}文件', f'*.{file_type.lower()}')]
        )
        
        if not save_path:
            return
            
        # 获取文件顺序
        ordered_files = []
        for i in range(self.order_list.size()):
            file_name = self.order_list.get(i)
            for file_path in self.files_to_convert:
                if os.path.basename(file_path) == file_name:
                    ordered_files.append(file_path)
                    break
        
        # 执行合并
        try:
            if file_type == 'PDF':
                from PyPDF2 import PdfMerger
                merger = PdfMerger()
                for file in ordered_files:
                    merger.append(file)
                merger.write(save_path)
                merger.close()
            else:  # Word
                from docx import Document
                from docxcompose.composer import Composer
                
                # 创建第一个文档
                doc = Document(ordered_files[0])
                composer = Composer(doc)
                
                # 添加其他文档
                for file in ordered_files[1:]:
                    doc_temp = Document(file)
                    composer.append(doc_temp)
                
                # 保存合并后的文档
                composer.save(save_path)
            
            # 添加到历史记录
            self.add_to_history(
                'merge',
                ordered_files[0],
                True,
                output_path=save_path,
                file_paths=ordered_files,
                file_type=file_type
            )
            
            # 关闭对话框
            dialog.destroy()
            
            # 询问是否打开合并后的文件
            if messagebox.askyesno("成功", "文件合并完成！是否打开合并后的文件？"):
                os.startfile(save_path)
                
        except Exception as e:
            messagebox.showerror("错误", f"合并文件时出错：{str(e)}")
            self.add_to_history(
                'merge',
                ordered_files[0],
                False,
                error=str(e),
                file_paths=ordered_files,
                file_type=file_type
            )

    def show_split_dialog(self):
        """显示文件分割对话框"""
        if not self.files_to_convert:
            messagebox.showwarning("警告", "请先选择要分割的文件！")
            return
            
        # 检查是否只选择了一个文件
        if len(self.files_to_convert) > 1:
            messagebox.showwarning("警告", "一次只能分割一个文件！")
            return
            
        file_path = self.files_to_convert[0]
        file_ext = os.path.splitext(file_path)[1].lower()
        
        # 检查文件类型
        if file_ext not in ['.pdf', '.docx', '.doc']:
            messagebox.showerror("错误", "只支持分割PDF和Word文件！")
            return
            
        # 创建分割对话框
        dialog = tk.Toplevel(self.root)
        dialog.title("文件分割")
        dialog.geometry("500x500")
        dialog.resizable(False, False)
        dialog.transient(self.root)
        dialog.grab_set()
        
        # 设置对话框居中显示
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f'{width}x{height}+{x}+{y}')
        
        # 创建分割选项框架
        split_frame = ttk.Frame(dialog, padding="20")
        split_frame.pack(fill=tk.BOTH, expand=True)
        
        # 文件信息
        ttk.Label(
            split_frame,
            text=f"文件名: {os.path.basename(file_path)}",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 10))
        
        # 分割方式
        ttk.Label(
            split_frame,
            text="分割方式:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.split_type_var = tk.StringVar(value="pages" if file_ext == '.pdf' else "sections")
        if file_ext == '.pdf':
            ttk.Radiobutton(
                split_frame,
                text="按页数分割",
                variable=self.split_type_var,
                value="pages",
                bootstyle="info-toolbutton"
            ).pack(anchor=tk.W, pady=2)
        else:
            ttk.Radiobutton(
                split_frame,
                text="按章节分割",
                variable=self.split_type_var,
                value="sections",
                bootstyle="info-toolbutton"
            ).pack(anchor=tk.W, pady=2)
        
        # 分割参数
        param_frame = ttk.Frame(split_frame)
        param_frame.pack(fill=tk.X, pady=10)
        
        if file_ext == '.pdf':
            ttk.Label(
                param_frame,
                text="每份页数:",
                font=('微软雅黑', 10),
                bootstyle="light"
            ).pack(side=tk.LEFT, padx=5)
            
            self.pages_per_split_var = tk.StringVar(value="10")
            ttk.Entry(
                param_frame,
                textvariable=self.pages_per_split_var,
                width=10
            ).pack(side=tk.LEFT, padx=5)
        else:
            ttk.Label(
                param_frame,
                text="分割标记:",
                font=('微软雅黑', 10),
                bootstyle="light"
            ).pack(side=tk.LEFT, padx=5)
            
            self.section_marker_var = tk.StringVar(value="###")
            ttk.Entry(
                param_frame,
                textvariable=self.section_marker_var,
                width=10
            ).pack(side=tk.LEFT, padx=5)
        
        # 输出目录
        ttk.Label(
            split_frame,
            text="输出目录:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.output_dir_var = tk.StringVar()
        ttk.Entry(
            split_frame,
            textvariable=self.output_dir_var,
            width=40
        ).pack(fill=tk.X, pady=(0, 10))
        
        ttk.Button(
            split_frame,
            text="选择目录",
            command=lambda: self.select_output_dir(self.output_dir_var),
            width=10,
            bootstyle="info"
        ).pack(anchor=tk.W, pady=(0, 10))
        
        # 创建按钮框架
        button_frame = ttk.Frame(split_frame)
        button_frame.pack(fill=tk.X, pady=20)
        
        # 分割按钮
        split_button = ttk.Button(
            button_frame,
            text="开始分割",
            command=lambda: self.start_split(file_path, dialog),
            width=10,
            bootstyle="success"
        )
        split_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 取消按钮
        cancel_button = ttk.Button(
            button_frame,
            text="取消",
            command=dialog.destroy,
            width=10,
            bootstyle="danger"
        )
        cancel_button.pack(side=tk.LEFT, padx=5, expand=True)
    
    def select_output_dir(self, var):
        """选择输出目录"""
        dir_path = filedialog.askdirectory()
        if dir_path:
            var.set(dir_path)
    
    def start_split(self, file_path, dialog):
        """开始分割文件"""
        output_dir = self.output_dir_var.get()
        if not output_dir:
            messagebox.showerror("错误", "请选择输出目录！")
            return
            
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            if file_ext == '.pdf':
                # PDF分割
                pages_per_split = int(self.pages_per_split_var.get())
                if pages_per_split <= 0:
                    raise ValueError("每份页数必须大于0")
                
                from PyPDF2 import PdfReader, PdfWriter
                reader = PdfReader(file_path)
                total_pages = len(reader.pages)
                
                output_files = []
                for i in range(0, total_pages, pages_per_split):
                    writer = PdfWriter()
                    end_page = min(i + pages_per_split, total_pages)
                    
                    for page_num in range(i, end_page):
                        writer.add_page(reader.pages[page_num])
                    
                    output_path = os.path.join(
                        output_dir,
                        f"{os.path.splitext(os.path.basename(file_path))[0]}_part{i//pages_per_split + 1}.pdf"
                    )
                    output_files.append(output_path)
                    
                    with open(output_path, 'wb') as output_file:
                        writer.write(output_file)
                
            else:
                # Word分割
                section_marker = self.section_marker_var.get()
                if not section_marker:
                    raise ValueError("分割标记不能为空")
                
                from docx import Document
                doc = Document(file_path)
                
                current_section = []
                section_num = 1
                output_files = []
                
                for para in doc.paragraphs:
                    if section_marker in para.text:
                        if current_section:
                            # 保存当前章节
                            new_doc = Document()
                            for p in current_section:
                                new_doc.add_paragraph(p.text, p.style)
                            
                            output_path = os.path.join(
                                output_dir,
                                f"{os.path.splitext(os.path.basename(file_path))[0]}_part{section_num}.docx"
                            )
                            output_files.append(output_path)
                            new_doc.save(output_path)
                            section_num += 1
                            current_section = []
                    else:
                        current_section.append(para)
                
                # 保存最后一个章节
                if current_section:
                    new_doc = Document()
                    for p in current_section:
                        new_doc.add_paragraph(p.text, p.style)
                    
                    output_path = os.path.join(
                        output_dir,
                        f"{os.path.splitext(os.path.basename(file_path))[0]}_part{section_num}.docx"
                    )
                    output_files.append(output_path)
                    new_doc.save(output_path)
            
            # 添加到历史记录
            self.add_to_history(
                'split',
                file_path,
                True,
                output_path=output_files[0] if output_files else None,
                output_files=output_files,
                split_type='pages' if file_ext == '.pdf' else 'sections',
                split_param=self.pages_per_split_var.get() if file_ext == '.pdf' else self.section_marker_var.get()
            )
            
            # 关闭对话框
            dialog.destroy()
            messagebox.showinfo("成功", "文件分割完成！")
            
        except ValueError as e:
            messagebox.showerror("错误", str(e))
        except Exception as e:
            messagebox.showerror("错误", f"分割文件时出错：{str(e)}")
            self.add_to_history(
                'split',
                file_path,
                False,
                error=str(e)
            )

    def show_conversion_settings(self):
        """显示转换参数设置对话框"""
        # 加载当前配置
        self.config = self.load_config()
        
        dialog = tk.Toplevel(self.root)
        dialog.title("转换参数设置")
        dialog.geometry("600x1150")
        dialog.resizable(False, False)
        dialog.transient(self.root)
        dialog.grab_set()
        
        # 设置对话框居中显示
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f'{width}x{height}+{x}+{y}')
        
        # 创建设置框架
        settings_frame = ttk.Frame(dialog, padding="20")
        settings_frame.pack(fill=tk.BOTH, expand=True)
        
        # 通用设置
        general_frame = ttk.Labelframe(
            settings_frame,
            text="通用设置",
            padding="10",
            bootstyle="info"
        )
        general_frame.pack(fill=tk.X, pady=10)
        
        # 自动打开选项
        ttk.Checkbutton(
            general_frame,
            text="转换完成后自动打开文件",
            variable=self.auto_open_var,
            bootstyle="info"
        ).pack(anchor=tk.W, pady=5)
        
        # PDF设置
        pdf_frame = ttk.Labelframe(
            settings_frame,
            text="PDF转换设置",
            padding="10",
            bootstyle="info"
        )
        pdf_frame.pack(fill=tk.X, pady=10)
        
        # PDF分辨率
        ttk.Label(
            pdf_frame,
            text="PDF分辨率 (DPI):",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.pdf_dpi_var = tk.StringVar(value=str(self.config.get('pdf_dpi', 300)))
        ttk.Entry(
            pdf_frame,
            textvariable=self.pdf_dpi_var,
            width=10
        ).pack(anchor=tk.W, pady=(0, 10))
        
        # PDF压缩质量
        ttk.Label(
            pdf_frame,
            text="PDF压缩质量:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.pdf_quality_var = tk.StringVar(value=self.config.get('pdf_quality', 'high'))
        ttk.Radiobutton(
            pdf_frame,
            text="高质量",
            variable=self.pdf_quality_var,
            value="high",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=2)
        
        ttk.Radiobutton(
            pdf_frame,
            text="标准",
            variable=self.pdf_quality_var,
            value="normal",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=2)
        
        ttk.Radiobutton(
            pdf_frame,
            text="低质量",
            variable=self.pdf_quality_var,
            value="low",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=2)
        
        # Word设置
        word_frame = ttk.Labelframe(
            settings_frame,
            text="Word转换设置",
            padding="10",
            bootstyle="info"
        )
        word_frame.pack(fill=tk.X, pady=10)
        
        # Word格式
        ttk.Label(
            word_frame,
            text="Word格式:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.word_format_var = tk.StringVar(value=self.config.get('word_format', 'docx'))
        ttk.Radiobutton(
            word_frame,
            text="DOCX (新版)",
            variable=self.word_format_var,
            value="docx",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=2)
        
        ttk.Radiobutton(
            word_frame,
            text="DOC (旧版)",
            variable=self.word_format_var,
            value="doc",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=2)
        
        # Excel设置
        excel_frame = ttk.Labelframe(
            settings_frame,
            text="Excel转换设置",
            padding="10",
            bootstyle="info"
        )
        excel_frame.pack(fill=tk.X, pady=10)
        
        # Excel格式
        ttk.Label(
            excel_frame,
            text="Excel格式:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.excel_format_var = tk.StringVar(value=self.config.get('excel_format', 'xlsx'))
        ttk.Radiobutton(
            excel_frame,
            text="XLSX (新版)",
            variable=self.excel_format_var,
            value="xlsx",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=2)
        
        ttk.Radiobutton(
            excel_frame,
            text="XLS (旧版)",
            variable=self.excel_format_var,
            value="xls",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=2)
        
        # PowerPoint设置
        ppt_frame = ttk.Labelframe(
            settings_frame,
            text="PowerPoint转换设置",
            padding="10",
            bootstyle="info"
        )
        ppt_frame.pack(fill=tk.X, pady=10)
        
        # PowerPoint格式
        ttk.Label(
            ppt_frame,
            text="PowerPoint格式:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.ppt_format_var = tk.StringVar(value=self.config.get('ppt_format', 'pptx'))
        ttk.Radiobutton(
            ppt_frame,
            text="PPTX (新版)",
            variable=self.ppt_format_var,
            value="pptx",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=2)
        
        ttk.Radiobutton(
            ppt_frame,
            text="PPT (旧版)",
            variable=self.ppt_format_var,
            value="ppt",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=2)
        
        # 创建按钮框架
        button_frame = ttk.Frame(settings_frame)
        button_frame.pack(fill=tk.X, pady=20)
        
        # 保存按钮
        save_button = ttk.Button(
            button_frame,
            text="保存",
            command=lambda: self.save_conversion_settings(dialog),
            width=10,
            bootstyle="success"
        )
        save_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 取消按钮
        cancel_button = ttk.Button(
            button_frame,
            text="取消",
            command=dialog.destroy,
            width=10,
            bootstyle="danger"
        )
        cancel_button.pack(side=tk.LEFT, padx=5, expand=True)
    
    def save_conversion_settings(self, dialog):
        """保存转换设置"""
        try:
            # 验证PDF DPI
            dpi = int(self.pdf_dpi_var.get())
            if dpi < 72 or dpi > 1200:
                raise ValueError("PDF分辨率必须在72-1200之间")
            
            # 更新配置
            self.config.update({
                'pdf_dpi': dpi,
                'pdf_quality': self.pdf_quality_var.get(),
                'word_format': self.word_format_var.get(),
                'excel_format': self.excel_format_var.get(),
                'ppt_format': self.ppt_format_var.get(),
                'auto_open': self.auto_open_var.get()
            })
            
            # 保存配置
            self.save_config()
            
            # 关闭对话框
            dialog.destroy()
            messagebox.showinfo("成功", "转换参数设置已保存！")
            
        except ValueError as e:
            messagebox.showerror("错误", str(e))
        except Exception as e:
            messagebox.showerror("错误", f"保存设置时出错：{str(e)}")

    def save_auto_open_setting(self):
        """保存自动打开文件设置"""
        self.config['auto_open'] = self.auto_open_var.get()
        self.save_config()

    def reconvert_history(self):
        """重新执行选中的历史记录"""
        selected_indices = self.history_list.curselection()
        if not selected_indices:
            messagebox.showwarning("警告", "请选择要重新执行的历史记录！")
            return
            
        selected_record = self.history[selected_indices[0]]
        if not selected_record['success']:
            messagebox.showwarning("警告", "选中的历史记录执行失败，无法重新执行！")
            return
            
        operation = selected_record['operation']
        
        if operation == 'convert':
            # 重新转换
            self.files_to_convert = [selected_record['file_path']]
            self.update_file_list()
            self.status_label.config(text=f"已选择历史文件: {os.path.basename(selected_record['file_path'])}")
            self.start_conversion()
            
        elif operation == 'split':
            # 重新分割
            self.files_to_convert = [selected_record['file_path']]
            self.update_file_list()
            
            # 创建分割对话框
            dialog = tk.Toplevel(self.root)
            dialog.title("文件分割")
            dialog.geometry("500x500")
            dialog.resizable(False, False)
            dialog.transient(self.root)
            dialog.grab_set()
            
            # 设置对话框居中显示
            dialog.update_idletasks()
            width = dialog.winfo_width()
            height = dialog.winfo_height()
            x = (dialog.winfo_screenwidth() // 2) - (width // 2)
            y = (dialog.winfo_screenheight() // 2) - (height // 2)
            dialog.geometry(f'{width}x{height}+{x}+{y}')
            
            # 创建分割选项框架
            split_frame = ttk.Frame(dialog, padding="20")
            split_frame.pack(fill=tk.BOTH, expand=True)
            
            # 文件信息
            ttk.Label(
                split_frame,
                text=f"文件名: {os.path.basename(selected_record['file_path'])}",
                font=('微软雅黑', 10),
                bootstyle="light"
            ).pack(anchor=tk.W, pady=(0, 10))
            
            # 分割方式
            ttk.Label(
                split_frame,
                text="分割方式:",
                font=('微软雅黑', 10),
                bootstyle="light"
            ).pack(anchor=tk.W, pady=(0, 5))
            
            self.split_type_var = tk.StringVar(value=selected_record['split_type'])
            if selected_record['split_type'] == 'pages':
                ttk.Radiobutton(
                    split_frame,
                    text="按页数分割",
                    variable=self.split_type_var,
                    value="pages",
                    bootstyle="info-toolbutton"
                ).pack(anchor=tk.W, pady=2)
            else:
                ttk.Radiobutton(
                    split_frame,
                    text="按章节分割",
                    variable=self.split_type_var,
                    value="sections",
                    bootstyle="info-toolbutton"
                ).pack(anchor=tk.W, pady=2)
            
            # 分割参数
            param_frame = ttk.Frame(split_frame)
            param_frame.pack(fill=tk.X, pady=10)
            
            if selected_record['split_type'] == 'pages':
                ttk.Label(
                    param_frame,
                    text="每份页数:",
                    font=('微软雅黑', 10),
                    bootstyle="light"
                ).pack(side=tk.LEFT, padx=5)
                
                self.pages_per_split_var = tk.StringVar(value=selected_record['split_param'])
                ttk.Entry(
                    param_frame,
                    textvariable=self.pages_per_split_var,
                    width=10
                ).pack(side=tk.LEFT, padx=5)
            else:
                ttk.Label(
                    param_frame,
                    text="分割标记:",
                    font=('微软雅黑', 10),
                    bootstyle="light"
                ).pack(side=tk.LEFT, padx=5)
                
                self.section_marker_var = tk.StringVar(value=selected_record['split_param'])
                ttk.Entry(
                    param_frame,
                    textvariable=self.section_marker_var,
                    width=10
                ).pack(side=tk.LEFT, padx=5)
            
            # 输出目录
            ttk.Label(
                split_frame,
                text="输出目录:",
                font=('微软雅黑', 10),
                bootstyle="light"
            ).pack(anchor=tk.W, pady=(0, 5))
            
            self.output_dir_var = tk.StringVar(value=os.path.dirname(selected_record['output_path']))
            ttk.Entry(
                split_frame,
                textvariable=self.output_dir_var,
                width=40
            ).pack(fill=tk.X, pady=(0, 10))
            
            ttk.Button(
                split_frame,
                text="选择目录",
                command=lambda: self.select_output_dir(self.output_dir_var),
                width=10,
                bootstyle="info"
            ).pack(anchor=tk.W, pady=(0, 10))
            
            # 创建按钮框架
            button_frame = ttk.Frame(split_frame)
            button_frame.pack(fill=tk.X, pady=20)
            
            # 分割按钮
            split_button = ttk.Button(
                button_frame,
                text="开始分割",
                command=lambda: self.start_split(selected_record['file_path'], dialog),
                width=10,
                bootstyle="success"
            )
            split_button.pack(side=tk.LEFT, padx=5, expand=True)
            
            # 取消按钮
            cancel_button = ttk.Button(
                button_frame,
                text="取消",
                command=dialog.destroy,
                width=10,
                bootstyle="danger"
            )
            cancel_button.pack(side=tk.LEFT, padx=5, expand=True)
            
        elif operation == 'merge':
            # 重新合并
            self.files_to_convert = selected_record['file_paths']
            self.update_file_list()
            
            # 创建合并对话框
            dialog = tk.Toplevel(self.root)
            dialog.title("文件合并")
            dialog.geometry("500x650")
            dialog.resizable(False, False)
            dialog.transient(self.root)
            dialog.grab_set()
            
            # 设置对话框居中显示
            dialog.update_idletasks()
            width = dialog.winfo_width()
            height = dialog.winfo_height()
            x = (dialog.winfo_screenwidth() // 2) - (width // 2)
            y = (dialog.winfo_screenheight() // 2) - (height // 2)
            dialog.geometry(f'{width}x{height}+{x}+{y}')
            
            # 创建合并选项框架
            merge_frame = ttk.Frame(dialog, padding="20")
            merge_frame.pack(fill=tk.BOTH, expand=True)
            
            # 文件类型标签
            file_type = selected_record['file_type']
            ttk.Label(
                merge_frame,
                text=f"合并 {file_type} 文件",
                font=('微软雅黑', 12, 'bold'),
                bootstyle="light"
            ).pack(anchor=tk.W, pady=(0, 20))
            
            # 输出文件名
            ttk.Label(
                merge_frame,
                text="输出文件名:",
                font=('微软雅黑', 10),
                bootstyle="light"
            ).pack(anchor=tk.W, pady=(0, 5))
            
            self.output_name_var = tk.StringVar(value=os.path.basename(selected_record['output_path']))
            ttk.Entry(
                merge_frame,
                textvariable=self.output_name_var,
                width=40
            ).pack(fill=tk.X, pady=(0, 20))
            
            # 文件顺序
            ttk.Label(
                merge_frame,
                text="文件顺序:",
                font=('微软雅黑', 10),
                bootstyle="light"
            ).pack(anchor=tk.W, pady=(0, 10))
            
            # 创建文件顺序列表
            self.order_list = tk.Listbox(
                merge_frame,
                width=50,
                height=8,
                font=('微软雅黑', 10),
                bg='#2b3e50',
                fg='white',
                selectbackground='#0078d7',
                selectforeground='white',
                relief=tk.FLAT
            )
            self.order_list.pack(fill=tk.BOTH, expand=True, pady=5)
            
            # 添加滚动条
            scrollbar = ttk.Scrollbar(
                merge_frame,
                orient=tk.VERTICAL,
                command=self.order_list.yview,
                bootstyle="round"
            )
            scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
            self.order_list.config(yscrollcommand=scrollbar.set)
            
            # 添加文件到顺序列表
            for file in selected_record['file_paths']:
                self.order_list.insert(tk.END, os.path.basename(file))
            
            # 创建按钮框架
            button_frame = ttk.Frame(merge_frame)
            button_frame.pack(fill=tk.X, pady=20)
            
            # 上移按钮
            up_button = ttk.Button(
                button_frame,
                text="上移",
                command=self.move_up,
                width=10,
                bootstyle="info"
            )
            up_button.pack(side=tk.LEFT, padx=5, expand=True)
            
            # 下移按钮
            down_button = ttk.Button(
                button_frame,
                text="下移",
                command=self.move_down,
                width=10,
                bootstyle="info"
            )
            down_button.pack(side=tk.LEFT, padx=5, expand=True)
            
            # 应用按钮
            apply_button = ttk.Button(
                button_frame,
                text="合并",
                command=lambda: self.apply_merge(dialog, file_type),
                width=10,
                bootstyle="success"
            )
            apply_button.pack(side=tk.LEFT, padx=5, expand=True)
            
            # 取消按钮
            cancel_button = ttk.Button(
                button_frame,
                text="取消",
                command=dialog.destroy,
                width=10,
                bootstyle="danger"
            )
            cancel_button.pack(side=tk.LEFT, padx=5, expand=True)
            
        else:
            messagebox.showwarning("警告", "不支持重新执行此类型的操作！")
    
    def clear_history(self):
        """清空历史记录"""
        if messagebox.askyesno("确认", "确定要清空所有历史记录吗？\n此操作不可恢复！", icon='warning'):
            try:
                self.history = []
                self.save_history()
                self.update_history_list()
                self.status_label.config(text="历史记录已清空")
                messagebox.showinfo("成功", "历史记录已清空！")
            except Exception as e:
                messagebox.showerror("错误", f"清空历史记录时出错：{str(e)}")

    def select_files(self):
        """选择文件并更新文件列表"""
        filetypes = [
            ("所有支持的文件", "*.pdf;*.docx;*.doc;*.xlsx;*.xls;*.pptx;*.ppt"),
            ("PDF文件", "*.pdf"),
            ("Word文件", "*.docx;*.doc"),
            ("Excel文件", "*.xlsx;*.xls"),
            ("PowerPoint文件", "*.pptx;*.ppt")
        ]
        
        files = filedialog.askopenfilenames(
            title="选择文件",
            filetypes=filetypes
        )
        
        if files:
            self.files_to_convert.extend(files)
            self.update_file_list()
            self.status_label.config(text=f"已选择 {len(files)} 个文件，共 {len(self.files_to_convert)} 个文件")
    
    def update_file_list(self):
        """更新文件列表显示"""
        self.file_list.delete(0, tk.END)
        for file in self.files_to_convert:
            self.file_list.insert(tk.END, os.path.basename(file))

    def show_advanced_settings(self):
        """显示高级设置对话框"""
        dialog = tk.Toplevel(self.root)
        dialog.title("高级设置")
        dialog.geometry("600x1350")
        dialog.resizable(False, False)
        dialog.transient(self.root)
        dialog.grab_set()
        
        # 设置对话框居中显示
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f'{width}x{height}+{x}+{y}')
        
        # 创建设置框架
        settings_frame = ttk.Frame(dialog, padding="20")
        settings_frame.pack(fill=tk.BOTH, expand=True)
        
        # 主题颜色设置
        theme_frame = ttk.Labelframe(
            settings_frame,
            text="主题颜色设置",
            padding="10",
            bootstyle="info"
        )
        theme_frame.pack(fill=tk.X, pady=10)
        
        # 主题选择
        ttk.Label(
            theme_frame,
            text="选择主题:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        # 创建主题选择框架
        theme_select_frame = ttk.Frame(theme_frame)
        theme_select_frame.pack(fill=tk.X, pady=(0, 10))
        
        valid_themes = ['darkly', 'solar', 'superhero', 'cyborg', 'vapor', 'cosmo', 'flatly', 'journal', 'litera', 'lumen', 'minty', 'pulse', 'sandstone', 'united', 'yeti']
        self.theme_var = tk.StringVar(value=self.config.get('theme', 'darkly'))
        
        # 使用Radiobuttons替代Combobox
        for i, theme in enumerate(valid_themes):
            ttk.Radiobutton(
                theme_select_frame,
                text=theme.capitalize(),
                variable=self.theme_var,
                value=theme,
                bootstyle="info-toolbutton"
            ).grid(row=i//3, column=i%3, padx=5, pady=2, sticky='w')

        # 自定义颜色
        ttk.Label(
            theme_frame,
            text="自定义颜色:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(10, 5))
        
        # 自定义颜色设置
        ttk.Label(
            theme_frame,
            text="前景色:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.fg_color_var = tk.StringVar(value=self.config.get('fg_color', '#ffffff'))
        fg_color_entry = ttk.Entry(
            theme_frame,
            textvariable=self.fg_color_var,
            width=10
        )
        fg_color_entry.pack(side=tk.LEFT, padx=5)
        
        ttk.Label(
            theme_frame,
            text="背景色:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.bg_color_var = tk.StringVar(value=self.config.get('bg_color', '#2b3e50'))
        bg_color_entry = ttk.Entry(
            theme_frame,
            textvariable=self.bg_color_var,
            width=10
        )
        bg_color_entry.pack(side=tk.LEFT, padx=5)
        
        # 快捷键设置
        shortcut_frame = ttk.Labelframe(
            settings_frame,
            text="快捷键设置",
            padding="10",
            bootstyle="info"
        )
        shortcut_frame.pack(fill=tk.X, pady=10)
        
        # 选择文件快捷键
        ttk.Label(
            shortcut_frame,
            text="选择文件快捷键:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.select_shortcut_var = tk.StringVar(value=self.config.get('select_shortcut', '<Control-o>'))
        ttk.Entry(
            shortcut_frame,
            textvariable=self.select_shortcut_var,
            width=20
        ).pack(anchor=tk.W, pady=(0, 10))
        
        # 开始转换快捷键
        ttk.Label(
            shortcut_frame,
            text="开始转换快捷键:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.convert_shortcut_var = tk.StringVar(value=self.config.get('convert_shortcut', '<Control-r>'))
        ttk.Entry(
            shortcut_frame,
            textvariable=self.convert_shortcut_var,
            width=20
        ).pack(anchor=tk.W, pady=(0, 10))
        
        # 添加快捷键说明
        ttk.Label(
            shortcut_frame,
            text="快捷键格式说明：\n使用 <Control- 表示 Ctrl 键\n例如：<Control-o> 表示 Ctrl+O\n注意：必须包含尖括号 <>",
            font=('微软雅黑', 9),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(10, 0))
        
        # 默认保存路径设置
        save_path_frame = ttk.Labelframe(
            settings_frame,
            text="默认保存路径设置",
            padding="10",
            bootstyle="info"
        )
        save_path_frame.pack(fill=tk.X, pady=10)
        
        # 默认保存路径
        ttk.Label(
            save_path_frame,
            text="默认保存路径:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.default_save_path_var = tk.StringVar(value=self.config.get('default_save_path', ''))
        ttk.Entry(
            save_path_frame,
            textvariable=self.default_save_path_var,
            width=50
        ).pack(fill=tk.X, pady=(0, 10))
        
        ttk.Button(
            save_path_frame,
            text="选择路径",
            command=lambda: self.select_default_save_path(self.default_save_path_var),
            width=10,
            bootstyle="info"
        ).pack(anchor=tk.W, pady=(0, 10))
        
        # 创建按钮框架
        button_frame = ttk.Frame(settings_frame)
        button_frame.pack(fill=tk.X, pady=20)
        
        # 保存按钮
        save_button = ttk.Button(
            button_frame,
            text="保存",
            command=lambda: self.save_advanced_settings(dialog),
            width=10,
            bootstyle="success"
        )
        save_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 取消按钮
        cancel_button = ttk.Button(
            button_frame,
            text="取消",
            command=dialog.destroy,
            width=10,
            bootstyle="danger"
        )
        cancel_button.pack(side=tk.LEFT, padx=5, expand=True)
    
    def select_default_save_path(self, var):
        """选择默认保存路径"""
        path = filedialog.askdirectory()
        if path:
            var.set(path)
    
    def save_advanced_settings(self, dialog):
        """保存高级设置"""
        try:
            old_theme = self.config.get('theme')
            new_theme = self.theme_var.get()
            
            # 更新配置
            self.config.update({
                'theme': new_theme,
                'fg_color': self.fg_color_var.get(),
                'bg_color': self.bg_color_var.get(),
                'select_shortcut': self.select_shortcut_var.get(),
                'convert_shortcut': self.convert_shortcut_var.get(),
                'default_save_path': self.default_save_path_var.get()
            })
            
            # 保存配置
            self.save_config()
            
            # 如果主题发生变化，重新创建style对象
            if old_theme != new_theme:
                self.style = ttk.Style()
                self.style.theme_use(new_theme)
            
            # 更新所有窗口的主题
            self.update_all_windows_theme()
            
            # 关闭对话框
            dialog.destroy()
            messagebox.showinfo("成功", "高级设置已保存！")
            
        except Exception as e:
            messagebox.showerror("错误", f"保存设置时出错：{str(e)}")
    
    def update_all_windows_theme(self):
        """更新所有窗口的主题"""
        try:
            # 根据主题设置合适的颜色
            theme = self.config.get('theme', 'darkly')
            is_dark_theme = theme in ['darkly', 'solar', 'superhero', 'cyborg', 'vapor']
            
            # 设置默认颜色
            if is_dark_theme:
                default_fg = '#ffffff'
                default_bg = '#2b3e50'
                button_bg = '#3498db'
                button_fg = '#ffffff'
            else:
                default_fg = '#212529'
                default_bg = '#ffffff'
                button_bg = '#0d6efd'
                button_fg = '#ffffff'
            
            # 更新基本样式
            self.style.configure('TFrame', background=default_bg)
            self.style.configure('TLabel', background=default_bg, foreground=default_fg)
            self.style.configure('TButton', background=button_bg, foreground=button_fg)
            self.style.configure('TEntry', fieldbackground=default_bg, foreground=default_fg)
            self.style.configure('TLabelframe', background=default_bg, foreground=default_fg)
            self.style.configure('TLabelframe.Label', background=default_bg, foreground=default_fg)
            self.style.configure('TRadiobutton', background=default_bg, foreground=default_fg)
            
            # 更新标题样式
            if hasattr(self, 'header_label'):
                self.header_label.configure(foreground=default_fg, background=default_bg)
            
            # 更新文件列表样式
            if hasattr(self, 'file_list'):
                self.file_list.configure(bg=default_bg, fg=default_fg,
                                       selectbackground=button_bg,
                                       selectforeground=button_fg)
            
            # 更新历史记录列表样式
            if hasattr(self, 'history_list'):
                self.history_list.configure(bg=default_bg, fg=default_fg,
                                         selectbackground=button_bg,
                                         selectforeground=button_fg)
            
            # 更新进度条样式
            if hasattr(self, 'progress'):
                self.style.configure('Horizontal.TProgressbar',
                                   background=button_bg,
                                   troughcolor=default_bg)
            
            # 更新状态标签样式
            if hasattr(self, 'status_label'):
                self.status_label.configure(foreground=default_fg, background=default_bg)
            
            # 更新统计标签样式
            if hasattr(self, 'stats_label'):
                self.stats_label.configure(foreground=default_fg, background=default_bg)
            
            # 更新所有按钮样式
            for widget in self.root.winfo_children():
                if isinstance(widget, ttk.Button):
                    widget.configure(style='TButton')
            
            # 强制更新主窗口
            self.root.update_idletasks()
            
        except Exception as e:
            print(f"更新主题时出现错误：{str(e)}")

    def bind_shortcuts(self):
        """绑定快捷键"""
        # 移除旧的绑定
        self.root.unbind('<Control-o>')
        self.root.unbind('<Control-O>')
        self.root.unbind('<Control-r>')
        self.root.unbind('<Control-R>')
        
        # 绑定新的快捷键
        select_shortcut = self.config.get('select_shortcut', '<Control-o>')
        convert_shortcut = self.config.get('convert_shortcut', '<Control-r>')
        
        # 确保快捷键格式正确
        if not select_shortcut.startswith('<'):
            select_shortcut = f'<{select_shortcut}'
        if not select_shortcut.endswith('>'):
            select_shortcut = f'{select_shortcut}>'
            
        if not convert_shortcut.startswith('<'):
            convert_shortcut = f'<{convert_shortcut}'
        if not convert_shortcut.endswith('>'):
            convert_shortcut = f'{convert_shortcut}>'
        
        if select_shortcut:
            self.root.bind(select_shortcut, lambda e: self.select_files())
        if convert_shortcut:
            self.root.bind(convert_shortcut, lambda e: self.start_conversion())

    def show_encrypt_dialog(self):
        """显示加密/解密对话框"""
        if not self.files_to_convert:
            messagebox.showwarning("警告", "请先选择要加密/解密的文件！")
            return
            
        # 检查是否只选择了一个文件
        if len(self.files_to_convert) > 1:
            messagebox.showwarning("警告", "一次只能加密/解密一个文件！")
            return
            
        file_path = self.files_to_convert[0]
        file_ext = os.path.splitext(file_path)[1].lower()
        
        # 检查文件类型
        if file_ext not in ['.pdf', '.docx', '.doc']:
            messagebox.showerror("错误", "只支持加密/解密PDF和Word文件！")
            return
            
        # 创建加密/解密对话框
        dialog = tk.Toplevel(self.root)
        dialog.title("文件加密/解密")
        dialog.geometry("500x510")
        dialog.resizable(False, False)
        dialog.transient(self.root)
        dialog.grab_set()
        
        # 设置对话框居中显示
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f'{width}x{height}+{x}+{y}')
        
        # 创建主框架
        main_frame = ttk.Frame(dialog, padding="20")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # 文件信息
        ttk.Label(
            main_frame,
            text=f"文件名: {os.path.basename(file_path)}",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 20))
        
        # 操作选择
        self.operation_var = tk.StringVar(value="encrypt")
        ttk.Radiobutton(
            main_frame,
            text="加密文件",
            variable=self.operation_var,
            value="encrypt",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=2)
        
        ttk.Radiobutton(
            main_frame,
            text="解密文件",
            variable=self.operation_var,
            value="decrypt",
            bootstyle="info-toolbutton"
        ).pack(anchor=tk.W, pady=2)
        
        # 密码输入
        ttk.Label(
            main_frame,
            text="密码:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(20, 5))
        
        self.password_var = tk.StringVar()
        password_entry = ttk.Entry(
            main_frame,
            textvariable=self.password_var,
            show="*",
            width=30
        )
        password_entry.pack(fill=tk.X, pady=(0, 10))
        
        # 确认密码（仅加密时显示）
        self.confirm_frame = ttk.Frame(main_frame)
        self.confirm_frame.pack(fill=tk.X, pady=(0, 20))
        
        ttk.Label(
            self.confirm_frame,
            text="确认密码:",
            font=('微软雅黑', 10),
            bootstyle="light"
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.confirm_var = tk.StringVar()
        confirm_entry = ttk.Entry(
            self.confirm_frame,
            textvariable=self.confirm_var,
            show="*",
            width=30
        )
        confirm_entry.pack(fill=tk.X)
        
        # 绑定操作选择变化事件
        def on_operation_change(*args):
            if self.operation_var.get() == "encrypt":
                self.confirm_frame.pack(fill=tk.X, pady=(0, 20))
            else:
                self.confirm_frame.pack_forget()
        
        self.operation_var.trace_add("write", on_operation_change)
        
        # 创建按钮框架
        button_frame = ttk.Frame(main_frame)
        button_frame.pack(fill=tk.X, pady=20)
        
        # 确定按钮
        confirm_button = ttk.Button(
            button_frame,
            text="确定",
            command=lambda: self.start_encryption(file_path, dialog),
            width=10,
            bootstyle="success"
        )
        confirm_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 取消按钮
        cancel_button = ttk.Button(
            button_frame,
            text="取消",
            command=dialog.destroy,
            width=10,
            bootstyle="danger"
        )
        cancel_button.pack(side=tk.LEFT, padx=5, expand=True)
    
    def start_encryption(self, file_path, dialog):
        """开始加密/解密文件"""
        import os  # 确保在函数开始时导入os模块
        operation = self.operation_var.get()
        password = self.password_var.get()
        
        if not password:
            messagebox.showerror("错误", "请输入密码！")
            return
        
        if operation == "encrypt" and password != self.confirm_var.get():
            messagebox.showerror("错误", "两次输入的密码不一致！")
            return
        
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            output_path = os.path.splitext(file_path)[0] + "_加密" + file_ext if operation == "encrypt" else os.path.splitext(file_path)[0] + "_解密" + file_ext
            
            if file_ext in ['.pdf']:
                from PyPDF2 import PdfReader, PdfWriter
                
                if operation == "encrypt":
                    # 加密PDF
                    reader = PdfReader(file_path)
                    writer = PdfWriter()
                    
                    # 复制所有页面
                    for page in reader.pages:
                        writer.add_page(page)
                    
                    # 添加加密
                    writer.encrypt(password)
                    
                    # 保存加密后的文件
                    with open(output_path, "wb") as output_file:
                        writer.write(output_file)
                else:
                    # 解密PDF
                    reader = PdfReader(file_path)
                    if reader.is_encrypted:
                        try:
                            reader.decrypt(password)
                        except:
                            messagebox.showerror("错误", "密码错误！")
                            return
                        
                        writer = PdfWriter()
                        for page in reader.pages:
                            writer.add_page(page)
                        
                        with open(output_path, "wb") as output_file:
                            writer.write(output_file)
                    else:
                        messagebox.showerror("错误", "该PDF文件未加密！")
                        return
                    
            elif file_ext in ['.docx', '.doc']:
                try:
                    from win32com import client
                    import win32com.client
                    
                    # 获取完整的绝对路径
                    abs_file_path = os.path.abspath(file_path)
                    abs_output_path = os.path.abspath(output_path)
                    
                    # 创建 Word 应用程序实例
                    word = win32com.client.Dispatch("Word.Application")
                    word.Visible = False
                    
                    try:
                        if operation == "encrypt":
                            # 打开文档
                            doc = word.Documents.Open(abs_file_path)
                            
                            # 设置密码保护
                            doc.Password = password
                            
                            # 保存加密后的文档
                            doc.SaveAs2(abs_output_path)
                            doc.Close()
                        else:
                            try:
                                # 尝试打开加密文档
                                doc = word.Documents.Open(abs_file_path, PasswordDocument=password)
                                
                                # 移除密码保护
                                doc.Password = ""
                                
                                # 保存解密后的文档
                                doc.SaveAs2(abs_output_path)
                                doc.Close()
                            except Exception as e:
                                messagebox.showerror("错误", "密码错误或文件未加密！")
                                if os.path.exists(abs_output_path):
                                    os.remove(abs_output_path)
                                raise e
                    finally:
                        # 确保关闭 Word 应用程序
                        try:
                            word.Quit()
                        except:
                            pass
                    
                except Exception as e:
                    messagebox.showerror("错误", f"处理文件时出错：{str(e)}")
                    if os.path.exists(output_path):
                        os.remove(output_path)
                    return
            
            # 添加到历史记录
            self.add_to_history(
                'encrypt' if operation == "encrypt" else 'decrypt',
                file_path,
                True,
                output_path=output_path
            )
            
            # 关闭对话框
            dialog.destroy()
            messagebox.showinfo("成功", f"文件{'加密' if operation == 'encrypt' else '解密'}完成！")
            
        except Exception as e:
            messagebox.showerror("错误", f"{'加密' if operation == 'encrypt' else '解密'}文件时出错：{str(e)}")
            self.add_to_history(
                'encrypt' if operation == "encrypt" else 'decrypt',
                file_path,
                False,
                error=str(e)
            )

    def load_templates(self):
        """加载转换模板"""
        try:
            with open('templates.json', 'r', encoding='utf-8') as f:
                return json.load(f)
        except FileNotFoundError:
            return []
        except Exception as e:
            messagebox.showerror("错误", f"加载模板时出错：{str(e)}")
            return []
    
    def save_templates(self):
        """保存转换模板"""
        try:
            with open('templates.json', 'w', encoding='utf-8') as f:
                json.dump(self.templates, f, ensure_ascii=False, indent=4)
        except Exception as e:
            messagebox.showerror("错误", f"保存模板时出错：{str(e)}")
    
    def show_template_dialog(self):
        """显示模板管理对话框"""
        dialog = tk.Toplevel(self.root)
        dialog.title("模板管理")
        dialog.geometry("800x1000")
        dialog.resizable(False, False)
        dialog.transient(self.root)
        dialog.grab_set()
        
        # 设置对话框居中显示
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f'{width}x{height}+{x}+{y}')
        
        # 创建主框架
        main_frame = ttk.Frame(dialog, padding="20")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # 创建模板列表框架
        list_frame = ttk.Labelframe(main_frame, text="已保存的模板", padding="10")
        list_frame.pack(fill=tk.BOTH, expand=True, pady=(0, 20))
        
        # 创建模板列表
        self.template_list = tk.Listbox(
            list_frame,
            width=50,
            height=10,
            font=('微软雅黑', 10),
            bg='#2b3e50',
            fg='white',
            selectbackground='#0078d7',
            selectforeground='white',
            relief=tk.FLAT
        )
        self.template_list.pack(fill=tk.BOTH, expand=True, pady=5)
        
        # 添加滚动条
        scrollbar = ttk.Scrollbar(
            list_frame,
            orient=tk.VERTICAL,
            command=self.template_list.yview,
            bootstyle="round"
        )
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.template_list.config(yscrollcommand=scrollbar.set)
        
        # 更新模板列表
        self.update_template_list()
        
        # 创建新模板框架
        new_template_frame = ttk.Labelframe(main_frame, text="新建模板", padding="10")
        new_template_frame.pack(fill=tk.X, pady=(0, 20))
        
        # 模板名称
        ttk.Label(
            new_template_frame,
            text="模板名称:",
            font=('微软雅黑', 10)
        ).pack(anchor=tk.W, pady=(0, 5))
        
        self.template_name_var = tk.StringVar()
        ttk.Entry(
            new_template_frame,
            textvariable=self.template_name_var,
            width=40
        ).pack(fill=tk.X, pady=(0, 10))
        
        # 转换设置
        settings_frame = ttk.Frame(new_template_frame)
        settings_frame.pack(fill=tk.X, pady=5)
        
        # PDF质量设置
        ttk.Label(
            settings_frame,
            text="PDF质量:",
            font=('微软雅黑', 10)
        ).pack(side=tk.LEFT, padx=(0, 5))
        
        self.template_quality_var = tk.StringVar(value="高质量")
        quality_combo = ttk.Combobox(
            settings_frame,
            textvariable=self.template_quality_var,
            values=["高质量", "标准", "低质量"],
            width=10,
            state="readonly"
        )
        quality_combo.pack(side=tk.LEFT, padx=5)
        
        # 输出格式设置
        formats_frame = ttk.Frame(new_template_frame)
        formats_frame.pack(fill=tk.X, pady=5)
        
        ttk.Label(
            formats_frame,
            text="输出格式:",
            font=('微软雅黑', 10)
        ).pack(side=tk.LEFT, padx=(0, 5))
        
        self.template_format_var = tk.StringVar(value="pdf")
        format_combo = ttk.Combobox(
            formats_frame,
            textvariable=self.template_format_var,
            values=["pdf", "docx", "doc", "xlsx", "xls", "pptx", "ppt"],
            width=10,
            state="readonly"
        )
        format_combo.pack(side=tk.LEFT, padx=5)
        
        # 按钮框架
        button_frame = ttk.Frame(main_frame)
        button_frame.pack(fill=tk.X, pady=10)
        
        # 保存模板按钮
        save_button = ttk.Button(
            button_frame,
            text="保存模板",
            command=lambda: self.save_template(dialog),
            width=15,
            bootstyle="success"
        )
        save_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 应用模板按钮
        apply_button = ttk.Button(
            button_frame,
            text="应用模板",
            command=lambda: self.apply_template(dialog),
            width=15,
            bootstyle="info"
        )
        apply_button.pack(side=tk.LEFT, padx=5, expand=True)
        
        # 删除模板按钮
        delete_button = ttk.Button(
            button_frame,
            text="删除模板",
            command=self.delete_template,
            width=15,
            bootstyle="danger"
        )
        delete_button.pack(side=tk.LEFT, padx=5, expand=True)
    
    def update_template_list(self):
        """更新模板列表显示"""
        self.template_list.delete(0, tk.END)
        for template in self.templates:
            quality_text = self.quality_map.get(template['quality'], template['quality'])
            self.template_list.insert(tk.END, f"{template['name']} ({template['format']} - {quality_text})")
    
    def save_template(self, dialog):
        """保存新模板"""
        name = self.template_name_var.get().strip()
        if not name:
            messagebox.showerror("错误", "请输入模板名称！")
            return
            
        # 检查是否存在同名模板
        for template in self.templates:
            if template['name'] == name:
                if not messagebox.askyesno("确认", "已存在同名模板，是否覆盖？"):
                    return
                self.templates.remove(template)
                break
        
        # 创建新模板
        template = {
            'name': name,
            'quality': self.quality_reverse_map[self.template_quality_var.get()],
            'format': self.template_format_var.get(),
            'timestamp': datetime.datetime.now().timestamp()
        }
        
        self.templates.append(template)
        self.save_templates()
        self.update_template_list()
        messagebox.showinfo("成功", "模板保存成功！")
        
        # 清空输入
        self.template_name_var.set("")
    
    def apply_template(self, dialog):
        """应用选中的模板"""
        selection = self.template_list.curselection()
        if not selection:
            messagebox.showwarning("警告", "请选择要应用的模板！")
            return
        
        template = self.templates[selection[0]]
        
        # 应用模板设置
        self.config['pdf_quality'] = template['quality']
        if template['format'].endswith('x'):
            self.config[f"{template['format'][:-1]}_format"] = template['format']
        else:
            self.config[f"{template['format']}_format"] = template['format']
        
        self.save_config()
        messagebox.showinfo("成功", "模板应用成功！")
        dialog.destroy()
    
    def delete_template(self):
        """删除选中的模板"""
        selection = self.template_list.curselection()
        if not selection:
            messagebox.showwarning("警告", "请选择要删除的模板！")
            return
        
        if messagebox.askyesno("确认", "确定要删除选中的模板吗？"):
            del self.templates[selection[0]]
            self.save_templates()
            self.update_template_list()
            messagebox.showinfo("成功", "模板删除成功！")

    def create_update_window(self):
        """创建更新进度窗口"""
        self.update_window = tk.Toplevel(self.root)
        self.update_window.title("更新进度")
        self.update_window.geometry("400x200")
        self.update_window.resizable(False, False)
        self.update_window.withdraw()  # 初始时隐藏
        
        # 创建进度条
        self.update_progress = ttk.Progressbar(
            self.update_window,
            length=300,
            mode='determinate'
        )
        self.update_progress.pack(pady=20)
        
        # 创建状态标签
        self.update_status = ttk.Label(
            self.update_window,
            text="正在检查更新...",
            font=('微软雅黑', 10)
        )
        self.update_status.pack(pady=10)
        
        # 创建按钮框架
        button_frame = ttk.Frame(self.update_window)
        button_frame.pack(pady=10)
        
        # 取消按钮
        self.cancel_button = ttk.Button(
            button_frame,
            text="取消",
            command=self.cancel_update,
            width=10
        )
        self.cancel_button.pack(side=tk.LEFT, padx=5)
        
        # 安装按钮
        self.install_button = ttk.Button(
            button_frame,
            text="安装",
            command=self.install_update,
            width=10,
            state='disabled'
        )
        self.install_button.pack(side=tk.LEFT, padx=5)
    
    def check_for_updates(self):
        """检查更新"""
        self.update_window.deiconify()
        self.update_status.config(text="正在检查更新...")
        self.update_progress['value'] = 0
        self.cancel_button.config(state='normal')
        self.install_button.config(state='disabled')
        
        def check():
            update_info = self.update_manager.check_update()
            if update_info['available']:
                self.update_status.config(
                    text=f"发现新版本 {update_info['version']}\n{update_info['description']}"
                )
                self.install_button.config(state='normal')
            else:
                self.update_status.config(text="当前已是最新版本")
                self.root.after(2000, self.update_window.withdraw)
        
        Thread(target=check).start()
    
    def update_download_progress(self, progress):
        """更新下载进度"""
        self.update_progress['value'] = progress
        self.update_status.config(text=f"正在下载更新: {progress:.1f}%")
    
    def install_update(self):
        """安装更新"""
        self.install_button.config(state='disabled')
        self.cancel_button.config(state='disabled')
        self.update_status.config(text="正在安装更新...")
        
        def install():
            if self.update_manager.install_update():
                self.update_status.config(text="更新安装完成，请重启程序")
                self.install_button.config(text="重启", command=self.restart_app)
                self.install_button.config(state='normal')
            else:
                self.update_status.config(text="更新安装失败")
                self.cancel_button.config(state='normal')
        
        Thread(target=install).start()
    
    def cancel_update(self):
        """取消更新"""
        self.update_window.withdraw()
    
    def restart_app(self):
        """重启程序"""
        python = sys.executable
        os.execl(python, python, *sys.argv)

    def get_conversion_type(self, file_path):
        """根据文件扩展名获取转换类型"""
        ext = os.path.splitext(file_path)[1].lower()
        if ext in ['.doc', '.docx']:
            return 'word2pdf'  # 修正：改为 word2pdf
        elif ext in ['.pdf']:
            # 暂时只支持转 Word，后续可扩展
            return 'pdf2word'  # 修正：改为 pdf2word
        elif ext in ['.xls', '.xlsx']:
            return 'excel2pdf' # 修正：改为 excel2pdf
        elif ext in ['.ppt', '.pptx']:
            return 'ppt2pdf'   # 修正：改为 ppt2pdf
        # 可以在这里添加更多支持的转换类型，例如 pdf2excel, pdf2ppt
        else:
            return None

if __name__ == "__main__":
    root = tk.Tk()
    app = FileConverterApp(root)
    root.mainloop() 